from __future__ import annotations

import shutil
import subprocess
from pathlib import Path
from typing import Any

import fitz
import pdfplumber
from docx import Document
from openpyxl import Workbook
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Inches
from xhtml2pdf import pisa

try:
    from pdf2docx import Converter
except ImportError:
    Converter = None  # type: ignore


def _find_soffice() -> str | None:
    candidates = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "soffice",
    ]
    import shutil as sh

    for c in candidates:
        p = sh.which(c) if c == "soffice" else (c if Path(c).exists() else None)
        if p:
            return c if c != "soffice" else p
    return None


def office_to_pdf_libreoffice(src: Path, out_dir: Path) -> Path | None:
    soffice = _find_soffice()
    if not soffice:
        return None
    out_dir.mkdir(parents=True, exist_ok=True)
    cmd = [
        soffice,
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(out_dir),
        str(src),
    ]
    try:
        subprocess.run(cmd, check=True, capture_output=True, timeout=120)
    except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired):
        return None
    expected = out_dir / (src.stem + ".pdf")
    return expected if expected.exists() else None


def word_to_pdf(path: Path, out: Path) -> None:
    r = office_to_pdf_libreoffice(path, out.parent)
    if r and r.exists():
        shutil.move(r, out)
        return
    docx = Document(path)
    pdf = fitz.open()
    page = pdf.new_page()
    text = "\n".join(p.text for p in docx.paragraphs)
    page.insert_text((72, 72), text[:8000], fontsize=11)
    pdf.save(out.as_posix())
    pdf.close()


def excel_to_pdf(path: Path, out: Path) -> None:
    r = office_to_pdf_libreoffice(path, out.parent)
    if r and r.exists():
        shutil.move(r, out)
        return
    wb = load_workbook(path, read_only=True, data_only=True)
    pdf = fitz.open()
    for sheet in wb.sheetnames[:5]:
        ws = wb[sheet]
        lines: list[str] = []
        for row in ws.iter_rows(max_row=60, max_col=20, values_only=True):
            lines.append(" | ".join("" if v is None else str(v) for v in row))
        page = pdf.new_page()
        text = f"{sheet}\n" + "\n".join(lines)
        page.insert_text((40, 72), text[:7000], fontsize=8)
    wb.close()
    pdf.save(out.as_posix())
    pdf.close()


def powerpoint_to_pdf(path: Path, out: Path) -> None:
    r = office_to_pdf_libreoffice(path, out.parent)
    if r and r.exists():
        shutil.move(r, out)
        return
    prs = Presentation(path)
    pdf = fitz.open()
    for slide in prs.slides[:30]:
        page = pdf.new_page(width=slide.slide_width, height=slide.slide_height)
        page.insert_text((40, 40), f"Slide (text export)", fontsize=12)
    if pdf.page_count == 0:
        pdf.new_page()
    pdf.save(out.as_posix())
    pdf.close()


def html_to_pdf(html_content: str, out: Path) -> None:
    out.parent.mkdir(parents=True, exist_ok=True)
    with open(out, "wb") as f:
        pisa.CreatePDF(html_content.encode("utf-8"), dest=f, encoding="utf-8")


def pdf_to_word(path: Path, out: Path) -> None:
    if Converter is None:
        raise RuntimeError("pdf2docx not available")
    cv = Converter(path.as_posix())
    cv.convert(out.as_posix())
    cv.close()


def pdf_to_excel_tables(path: Path, out: Path) -> dict[str, Any]:
    rows: list[list[Any]] = []
    with pdfplumber.open(path) as pdf:
        for p in pdf.pages[:50]:
            tables = p.extract_tables() or []
            for tb in tables:
                for row in tb:
                    rows.append(row)
    if not rows:
        with pdfplumber.open(path) as pdf0:
            t0 = pdf0.pages[0].extract_text() if pdf0.pages else ""
        rows = [["Extracted text"], [t0 or ""]]
    xwb = Workbook()
    ws = xwb.active
    ws.title = "Extracted"
    for ri, row in enumerate(rows, start=1):
        for ci, cell in enumerate(row, start=1):
            ws.cell(row=ri, column=ci, value=cell)
    xwb.save(out)
    return {"rows": len(rows)}


def pdf_to_powerpoint(path: Path, out: Path) -> None:
    doc = fitz.open(path)
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for i in range(min(doc.page_count, 40)):
        slide = prs.slides.add_slide(blank)
        pix = doc[i].get_pixmap(matrix=fitz.Matrix(1.5, 1.5))
        img_path = out.parent / f"_slide_{i}.png"
        pix.save(img_path.as_posix())
        slide.shapes.add_picture(str(img_path), Inches(0.2), Inches(0.2), width=Inches(9))
        try:
            img_path.unlink()
        except OSError:
            pass
    prs.save(out.as_posix())
    doc.close()
